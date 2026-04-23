"""
backend/services/slides_renderer.py
──────────────────────────────────────
Slides renderer for artifact system.

Spec: L – Artifact System
Branch: engineering/master-list-closeout

Strategy (in order of availability):
  1. python-pptx — native PPTX generation (pip install python-pptx)
  2. HTML slides — Reveal.js-compatible standalone HTML (always available)

Takes:
  content: str (markdown) | dict {"slides": [{"title":..., "bullets":[...]}]}
  title:   str
  output_path: str (should end in .pptx or .html)

Returns the file path of the generated output.
"""

from __future__ import annotations

import json
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

logger = logging.getLogger(__name__)

ARTIFACT_STORAGE_DIR = os.environ.get("ARTIFACT_STORAGE_DIR", "/tmp/crucibai_artifacts")


def _parse_content(content: Union[str, Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Convert content to list of slide dicts: [{"title": ..., "bullets": [...]}]."""
    if isinstance(content, dict):
        if "slides" in content:
            return content["slides"]
        # Single slide
        return [{"title": content.get("title", "Slide"), "bullets": content.get("bullets", [])}]
    # Markdown: headings become slide titles, bullets become bullet points
    slides = []
    current: Optional[Dict[str, Any]] = None
    for line in str(content).split("\n"):
        stripped = line.strip()
        if stripped.startswith("# ") or stripped.startswith("## "):
            if current:
                slides.append(current)
            current = {"title": stripped.lstrip("# "), "bullets": []}
        elif stripped.startswith("- ") and current:
            current["bullets"].append(stripped[2:])
        elif stripped and current:
            current["bullets"].append(stripped)
    if current:
        slides.append(current)
    if not slides:
        slides = [{"title": "Presentation", "bullets": [str(content)[:200]]}]
    return slides


class SlidesRenderer:
    """Render content to PPTX or HTML slides."""

    async def render(
        self,
        *,
        content: Union[str, Dict[str, Any]],
        title: str,
        output_path: Optional[str] = None,
        format: str = "auto",   # pptx | html | auto
    ) -> str:
        os.makedirs(ARTIFACT_STORAGE_DIR, exist_ok=True)
        if output_path is None:
            import uuid
            output_path = os.path.join(ARTIFACT_STORAGE_DIR, f"{uuid.uuid4()}.pptx")

        slides = _parse_content(content)

        # Try PPTX first unless html forced
        if format in ("pptx", "auto"):
            result = self._try_pptx(slides, title, output_path)
            if result:
                return result

        # HTML fallback
        html_path = output_path.replace(".pptx", ".html")
        return self._render_html(slides, title, html_path)

    def _try_pptx(self, slides: List[Dict], title: str, output_path: str) -> Optional[str]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Title slide
            layout_title = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout_title)
            slide.shapes.title.text = title
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = f"CrucibAI · {len(slides)} slides"

            # Content slides
            layout_content = prs.slide_layouts[1]
            for s in slides:
                slide = prs.slides.add_slide(layout_content)
                slide.shapes.title.text = s.get("title", "")[:100]
                if slide.placeholders and len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for bullet in (s.get("bullets") or [])[:10]:
                        p = tf.add_paragraph()
                        p.text = str(bullet)[:200]
                        p.level = 0

            prs.save(output_path)
            logger.info("[SlidesRenderer] pptx: %s", output_path)
            return output_path
        except ImportError:
            return None
        except Exception as exc:
            logger.warning("[SlidesRenderer] pptx error: %s", exc)
            return None

    def _render_html(self, slides: List[Dict], title: str, output_path: str) -> str:
        import html as _html

        def esc(s: Any) -> str:
            return _html.escape(str(s))

        sections = []
        # Title slide
        sections.append(f"""<section>
  <h1>{esc(title)}</h1>
  <p style="color:#888">CrucibAI · {len(slides)} slides</p>
</section>""")

        for s in slides:
            bullets_html = "".join(f"<li>{esc(b)}</li>" for b in (s.get("bullets") or []))
            sections.append(f"""<section>
  <h2>{esc(s.get("title", ""))}</h2>
  <ul style="font-size:1.2em;text-align:left">{bullets_html}</ul>
</section>""")

        body = "\n".join(sections)
        html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{esc(title)}</title>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.0/reset.min.css">
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.0/reveal.min.css">
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.0/theme/black.min.css">
</head>
<body>
<div class="reveal">
  <div class="slides">
{body}
  </div>
</div>
<script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.0/reveal.min.js"></script>
<script>Reveal.initialize({{hash:true,transition:"slide"}});</script>
</body>
</html>"""
        Path(output_path).write_text(html_content, encoding="utf-8")
        logger.info("[SlidesRenderer] html: %s", output_path)
        return output_path


# Module-level singleton
slides_renderer = SlidesRenderer()
